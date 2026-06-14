# -*- coding: utf-8 -*-
"""
Generate arayeh-kamel.pptx — full 30-min presentation
Pastel pink/purple theme · CFG images · all mutants & kills
"""
from pathlib import Path
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))
from cfg_png_export import export_all  # noqa: E402

OUT = ROOT / "cognitive-frontend" / "public" / "arayeh-kamel.pptx"
OUT2 = ROOT / "silver_project" / "arayeh-kamel.pptx"
ASSETS = ROOT / "cognitive-frontend" / "public" / "pptx_assets"

# Theme
C_BG = RGBColor(0xFD, 0xF2, 0xF8)
C_PINK = RGBColor(0xEC, 0x48, 0x99)
C_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
C_TEXT = RGBColor(0x58, 0x1C, 0x87)
C_MUTED = RGBColor(0x93, 0x33, 0xEA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_PINK_SOFT = RGBColor(0xFB, 0xCF, 0xE8)
C_PURPLE_SOFT = RGBColor(0xE9, 0xD5, 0xFF)
C_GREEN = RGBColor(0x16, 0xA3, 0x4A)
C_RED = RGBColor(0xDC, 0x26, 0x26)

MUTANTS = [
    ("M-AOR-BS-01", "bubble_sort", "n-i-1 → n+i-1", "KILLED", "IndexError", "test_m01_bs_killed_index_error", "mutation"),
    ("M-AOR-BS-02", "bubble_sort", "n-i-1 → n-i+1", "KILLED", "IndexError پاس اول", "test_m02_bs_killed_index_error_first_pass", "mutation"),
    ("M-AOR-BS-03", "bubble_sort", "n-i-1 → n-i-2", "KILLED", "خروجی [b,c,d,a]≠sorted", "test_m03_bs_killed_wrong_output", "mutation"),
    ("M-AOR-BS-04", "bubble_sort", "n-i-1 → n*i-1", "KILLED", "بدون مرتب‌سازی", "test_m04_bs_killed_no_sorting", "mutation"),
    ("M-AOR-MS-01", "merge_sort", "n//2 → n//1", "KILLED", "RecursionError", "test_m01_ms_killed_recursion_error", "mutation"),
    ("M-AOR-MS-02", "merge_sort", "n//2 → n//3", "KILLED", "RecursionError n=2", "test_m02_ms_killed_recursion_error", "mutation"),
    ("M-AOR-MS-03", "merge_sort", "n//2 → (n+1)//2", "EQUIVALENT", "خروجی همیشه یکسان", "—", "—"),
    ("M-AOR-BN-01", "binary_search", "(L+R)//2 → (L-R)//2", "KILLED", "اندیس غلط / -1", "test_m01_bn_killed", "mutation"),
    ("M-AOR-BN-02", "binary_search", "//2 → *2", "KILLED", "IndexError", "test_m02_bn_killed_index_error", "mutation"),
    ("M-AOR-BN-03", "binary_search", "mid سقفی", "EQUIVALENT", "هر دو mid معتبر", "—", "—"),
    ("M-AOR-BN-04", "binary_search", "left=mid-1", "KILLED", "حلقه بی‌نهایت", "test_m04_bn timeout", "mutation"),
    ("M-AOR-BN-05", "binary_search", "right=mid+1", "KILLED", "حلقه بی‌نهایت", "test_m05_bn timeout", "mutation"),
    ("M-AOR-PC-01", "catalog", "total_before+1", "KILLED*", "meta غلط", "ACOC + assert total_before", "ACOC"),
    ("M-AOR-PC-02", "catalog", "total_after-1", "KILLED", "meta≠len(items)", "test_acoc_catalog_combinations", "ACOC"),
]


class Deck:
    def __init__(self):
        export_all()
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def slide(self):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        self._bg(s)
        self._header_bar(s)
        self._footer(s)
        return s

    def _bg(self, slide):
        f = slide.background.fill
        f.solid()
        f.fore_color.rgb = C_BG

    def _header_bar(self, slide):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_PINK
        bar.line.fill.background()
        bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.12), Inches(10), Inches(0.06))
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = C_PURPLE
        bar2.line.fill.background()

    def _footer(self, slide):
        foot = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(9.4), Inches(0.3))
        p = foot.text_frame.paragraphs[0]
        p.text = f"ctlg · تست نرم‌افزار · اسلاید {self.n}"
        p.font.size = Pt(9)
        p.font.color.rgb = C_MUTED
        p.alignment = PP_ALIGN.CENTER

    def title(self, slide, text, sub=None, y=0.35):
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(y), Inches(9.1), Inches(0.75))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = C_PINK
        p.alignment = PP_ALIGN.RIGHT
        if sub:
            tb2 = slide.shapes.add_textbox(Inches(0.45), Inches(y + 0.72), Inches(9.1), Inches(0.4))
            p2 = tb2.text_frame.paragraphs[0]
            p2.text = sub
            p2.font.size = Pt(13)
            p2.font.color.rgb = C_MUTED
            p2.alignment = PP_ALIGN.RIGHT

    def body(self, slide, lines, top=1.35, sz=14):
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(top), Inches(9.1), Inches(5.6))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(sz if not line.startswith("  ") else sz - 1)
            p.font.color.rgb = C_TEXT
            p.alignment = PP_ALIGN.RIGHT
            p.space_after = Pt(4)
            if line.strip().startswith(("assert", "def ", "for ", "if ", "search", "sort")):
                p.font.name = "Consolas"

    def table(self, slide, headers, rows, top=1.35, col_w=None):
        nr, nc = len(rows) + 1, len(headers)
        h = min(Inches(0.32 * nr + 0.3), Inches(5.5))
        shape = slide.shapes.add_table(nr, nc, Inches(0.4), Inches(top), Inches(9.2), h)
        tbl = shape.table
        if col_w:
            for j, w in enumerate(col_w):
                tbl.columns[j].width = Inches(w)
        for j, htxt in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = htxt
            c.fill.solid()
            c.fill.fore_color.rgb = C_PINK_SOFT
            for p in c.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(10)
                p.font.color.rgb = C_PURPLE
                p.alignment = PP_ALIGN.RIGHT
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                c = tbl.cell(i, j)
                c.text = str(val)
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.alignment = PP_ALIGN.RIGHT
                    if j == 3 and "KILLED" in str(val):
                        p.font.color.rgb = C_GREEN
                        p.font.bold = True
                    elif j == 3 and "EQUIVAL" in str(val):
                        p.font.color.rgb = C_PURPLE
                    elif j == 3 and "LIVE" in str(val):
                        p.font.color.rgb = C_RED
                    else:
                        p.font.color.rgb = C_TEXT

    def image(self, slide, name, left, top, width):
        path = ASSETS / name
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))

    def two_col(self, slide, left_lines, right_img, img_name, img_w=4.2):
        self.body(slide, left_lines, top=1.35, sz=13)
        self.image(slide, img_name, 0.45, 1.35, img_w)

    def build_title(self):
        s = self.slide()
        deco = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.5), Inches(7.6), Inches(4.5))
        deco.fill.solid()
        deco.fill.fore_color.rgb = C_WHITE
        deco.line.color.rgb = C_PINK_SOFT
        deco.line.width = Pt(2)
        tb = s.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(7.0), Inches(3.8))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lines = [
            ("تست نرم‌افزار", 32, C_PINK, True),
            ("پروژه کاتالوگ آزمون‌های شناختی", 20, C_PURPLE, True),
            ("", 8, C_TEXT, False),
            ("ACOC · CFG · Mutation Testing", 16, C_MUTED, False),
            ("UUT: process_catalog", 14, C_TEXT, False),
            ("", 8, C_TEXT, False),
            ("376 pytest  |  91.7% Mutation Score  |  14 AOR", 15, C_PINK, True),
        ]
        for i, (txt, sz, col, bold) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = txt
            p.font.size = Pt(sz)
            p.font.color.rgb = col
            p.font.bold = bold
            p.alignment = PP_ALIGN.CENTER

    def build_all(self):
        self.build_title()

        s = self.slide()
        self.title(s, "فهرست ارائه (~۳۰ دقیقه)", "۲۸ اسلاید — ACOC + CFG + Mutation + گراف‌ها")
        self.body(s, [
            "۱–۵   سناریو، UUT، الگوریتم‌ها در سامانه، جریان catalog",
            "۶–۱۰  ISP/ACOC — ۳۳۰ تست، فرمول، oracle، نمونه کد",
            "۱۱–۱۶ CFG — Node/Edge/Prime + گراف‌های bubble/linear/catalog",
            "۱۷–۲۳ Mutation — ۱۴ AOR، جدول kill، نمونه کد، PC-01",
            "۲۴–۲۷ نمودار score · پاسخ ۳ سؤال تکلیف",
            "۲۸    جمع‌بندی و منابع",
        ], sz=15)

        s = self.slide()
        self.title(s, "سناریو — سامانه سنجش شناختی")
        self.body(s, [
            "دانشجو/مدرس لیست آزمون‌های شناختی را می‌بیند.",
            "می‌تواند جستجو (query) و مرتب‌سازی (field، reverse) انجام دهد.",
            "",
            "جریان داده:",
            "  PostgreSQL → Django → catalog_bridge → process_catalog → JSON → React",
            "",
            "کاربر نام الگوریتم (bubble/merge/linear) نمی‌بیند.",
            "انتخاب خودکار در catalog_bridge.py انجام می‌شود.",
        ])

        s = self.slide()
        self.title(s, "معماری و UUT")
        self.table(s, ["لایه", "مسیر", "نقش"], [
            ("UUT اصلی", "silver_project/algorithms/catalog.py", "process_catalog"),
            ("مرتب‌سازی", "sorting.py", "bubble_sort, merge_sort"),
            ("جستجو", "searching.py", "linear_search, binary_search"),
            ("پل Django", "coglearning/assessment/catalog_bridge.py", "انتخاب algo + API"),
            ("UI", "cognitive-frontend/", "React — لیست آزمون"),
        ])

        s = self.slide()
        self.title(s, "الگوریتم‌ها — کاربرد واقعی")
        self.table(s, ["الگوریتم", "عملکرد", "زمان استفاده"], [
            ("linear_search", "زیررشته در title و description", "همیشه — production"),
            ("bubble_sort", "مرتب‌سازی مقایسه‌ای O(n²)", "لیست ≤ ۱۲ آیتم"),
            ("merge_sort", "تقسیم و حل O(n log n)", "لیست > ۱۲ آیتم"),
            ("binary_search", "تطبیق دقیق روی sorted", "فقط تست ACOC/CFG"),
        ], top=1.3)
        self.body(s, ["  search_algo = 'linear'  # همیشه", "  sort_algo = 'bubble' if count <= 12 else 'merge'"], top=4.2, sz=12)

        s = self.slide()
        self.title(s, "جریان process_catalog", "۵ مرحله — گراف در اسلاید بعد")
        self.body(s, [
            "PC1 — اعتبارسنجی: sort_field، sort_algo، search_algo (invalid → default)",
            "PC2 — filtered = _apply_search(items, query, search_algo, sort_field)",
            "  AS1: query خالی؟ → return همه",
            "  AS3: binary؟ → bubble_sort + binary_search : linear_search",
            "PC3 — sorted_items = _apply_sort(filtered, sort_algo, field, reverse)",
            "PC4 — return { items: sorted_items, meta: {...} }",
        ], sz=13)

        s = self.slide()
        self.title(s, "گراف CFG — process_catalog (UUT)")
        self.image(s, "cfg_process_catalog.png", 0.5, 1.25, 9.0)

        s = self.slide()
        self.title(s, "چرا ACOC؟ — Input Space Partitioning")
        self.body(s, [
            "تست همه حالات ورودی ممکن نیست (اندازه، ترتیب، query، ...).",
            "",
            "ISP: فضای ورودی را به بلوک‌های معنی‌دار تقسیم می‌کنیم.",
            "ACOC: All Combinations — ترکیب همه بلوک‌ها = پوشش systematic.",
            "",
            "۶ پارامتر × مقادیر گسسته = 324 ترکیب",
            "+ ۶ oracle مرتب‌سازی = 330 تست",
            "",
            "فایل: test_coverage.py → test_acoc_catalog_combinations",
        ])

        s = self.slide()
        self.title(s, "جدول پارامترهای ACOC")
        self.table(s, ["پارامتر", "بلوک‌ها", "n", "توضیح"], [
            ("size", "[], 1 item, 2+", "3", "خالی / تک / چند"),
            ("order", "sorted, reverse, random", "3", "ترتیب اولیه"),
            ("sort_field", "title, min_level, time_limit", "3", "کلید مرتب"),
            ("sort_algo", "bubble, merge", "2", "الگوریتم sort"),
            ("query", '""، آزمون، ناموجود', "3", "جستجو"),
            ("search_algo", "linear, binary", "2", "مسیر _apply_search"),
        ], col_w=[1.8, 3.5, 0.6, 2.5])

        s = self.slide()
        self.title(s, "فرمول تعداد تست ACOC")
        self.image(s, "test_overview_chart.png", 0.4, 1.3, 5.0)
        self.body(s, [
            "3×3×3×2×3×2 = 324",
            "",
            "+ 6 oracle:",
            "  2 sort_algo × 3 sort_field",
            "  → بررسی مرتب بودن خروجی",
            "",
            "جمع ACOC = 330",
            "CFG = 29  |  Mutation tests = 17",
            "کل pytest = 376 passed",
        ], top=1.35, sz=13)

        s = self.slide()
        self.title(s, "Oracle — مرجع صحت ACOC")
        self.body(s, [
            "assert 'items' in result and 'meta' in result",
            "",
            "assert result['meta']['total_after'] == len(result['items'])",
            "  → PC-02 mutant را kill می‌کند",
            "",
            "assert result['meta']['total_after'] <= result['meta']['total_before']",
            "",
            "assert result['meta']['total_before'] == len(items)",
            "  → PC-01 mutant را kill می‌کند (تقویت mutation score)",
            "",
            "برای sort: get_item_value(i) <= get_item_value(i+1)",
        ], sz=13)

        s = self.slide()
        self.title(s, "CFG — Control Flow Graph (فصل ۷)")
        self.body(s, [
            "کد → گراف: گره = statement/شرط · یال = جریان (+ T/F)",
            "",
            "Node Coverage (14 تست):",
            "  هر گره CFG حداقل یک‌بار اجرا شود",
            "",
            "Edge Coverage (9 تست):",
            "  هر یال و هر شاخه True/False پوشش داده شود",
            "",
            "Prime Path Coverage (6 تست):",
            "  مسیرهای اصلی — حلقه تو در تو، recursion، pipeline کامل",
        ], sz=14)

        s = self.slide()
        self.title(s, "گراف CFG — bubble_sort")
        self.image(s, "cfg_bubble_sort.png", 0.3, 1.2, 5.5)
        self.body(s, [
            "BS2: n<=1 → return",
            "BS4: حلقه بیرونی i",
            "BS6: حلقه داخلی j",
            "BS7: swap؟",
            "BS9: early exit",
            "",
            "Node: empty list → BS2=T",
            "Edge: [b,a] → BS7=T",
            "Prime: [d,c,b,a] → full loops",
        ], top=1.3, sz=12)

        s = self.slide()
        self.title(s, "گراف CFG — linear_search")
        self.image(s, "cfg_linear_search.png", 0.3, 1.2, 5.5)
        self.body(s, [
            "LS1: query خالی → همه index",
            "LS4: حلقه i روی items",
            "LS5: حلقه field",
            "LS6: query in value?",
            "LS7: append + break",
            "",
            "production: همیشه این مسیر",
            "جستجوی کاربر در UI",
        ], top=1.3, sz=12)

        s = self.slide()
        self.title(s, "نمونه تست‌های CFG")
        self.table(s, ["نوع", "نام تست", "ورودی", "پوشش"], [
            ("Node", "test_bubble_empty_list", "[]", "BS2→BS3"),
            ("Node", "test_merge_multiple", "[c,a,b]", "MS divide"),
            ("Edge", "test_bubble_swap_taken", "[b,a]", "BS7=T"),
            ("Edge", "test_catalog_binary_path", "binary", "AS3=T"),
            ("Prime", "test_bubble_full_loop", "[d,c,b,a]", "BS4→BS6"),
            ("Prime", "test_catalog_pipeline", "query+sort", "PC2→PC3"),
        ])

        s = self.slide()
        self.title(s, "Mutation Testing — AOR (فصل ۹)")
        self.body(s, [
            "Arithmetic Operator Replacement:",
            "  یک عملگر حسابی در سورس عوض می‌شود → Mutant",
            "",
            "اجرای همان تست روی mutant:",
            "  KILLED     — تست fail (AssertionError, Exception, خروجی غلط)",
            "  EQUIVALENT — همیشه خروجی = برنامه اصل (قابل kill نیست)",
            "  LIVE       — تست pass — suite ضعیف",
            "",
            "Mutation Score = Killed / (Total − Equivalent)",
            "  = 11 / (14 − 2) = 91.7%",
        ], sz=14)

        s = self.slide()
        self.title(s, "نمودار Mutation Score")
        self.image(s, "mutation_score_chart.png", 1.5, 1.4, 7.0)

        s = self.slide()
        self.title(s, "جدول کامل ۱۴ Mutant — وضعیت Kill")
        self.table(s, ["ID", "تابع", "تغییر AOR", "وضعیت", "اثر", "تست kill", "Suite"], MUTANTS,
                     top=1.25, col_w=[1.1, 1.0, 1.5, 0.9, 1.3, 1.8, 0.8])

        s = self.slide()
        self.title(s, "Mutantهای bubble_sort — ۴/۴ Killed")
        self.table(s, ["ID", "تغییر", "اثر", "تست"], [
            ("BS-01", "n+i-1", "IndexError", "pytest.raises(IndexError)"),
            ("BS-02", "n-i+1", "arr[n] OOB", "pytest.raises(IndexError)"),
            ("BS-03", "n-i-2", "[b,c,d,a]", "assert got != correct"),
            ("BS-04", "n*i-1", "range خالی", "assert still [b,a]"),
        ])
        self.body(s, ["فایل: mutants/sorting_mutants.py · test_mutation_killing.py"], top=4.5, sz=11)

        s = self.slide()
        self.title(s, "Mutantهای merge_sort")
        self.table(s, ["ID", "تغییر", "وضعیت", "توضیح"], [
            ("MS-01", "n//1 = n", "KILLED", "RecursionError — mid=n"),
            ("MS-02", "n//3", "KILLED", "RecursionError — mid=0 for n=2"),
            ("MS-03", "(n+1)//2", "EQUIVALENT", "تقسیم سقفی — sort یکسان"),
        ])
        self.body(s, [
            "merge در production: catalog > 12 آیتم",
            "Score merge (killable): 2/2 = 100%",
        ], top=3.8, sz=13)

        s = self.slide()
        self.title(s, "Mutantهای binary_search")
        self.table(s, ["ID", "تغییر", "وضعیت", "مکانیزم kill"], [
            ("BN-01", "(L-R)//2", "KILLED", "result != 2 for 'c'"),
            ("BN-02", "(L+R)*2", "KILLED", "IndexError"),
            ("BN-03", "(L+R+1)//2", "EQUIVALENT", "mid سقفی معتبر"),
            ("BN-04", "left=mid-1", "KILLED", "thread timeout — infinite loop"),
            ("BN-05", "right=mid+1", "KILLED", "thread timeout — infinite loop"),
        ])

        s = self.slide()
        self.title(s, "Mutantهای catalog + داستان PC-01")
        self.body(s, [
            "PC-02: total_after = len(sorted)-1",
            "  → kill توسط ACOC: total_after == len(items)",
            "",
            "PC-01: total_before = len(items)+1",
            "  قبل: فقط total_after <= total_before → LIVE",
            "  Score: 10/12 = 83.3%",
            "",
            "بعد از fix:",
            "  assert total_before == len(items)",
            "  → KILLED · Score: 11/12 = 91.7%",
            "",
            "درس: ACOC تعداد را پوشش می‌دهد — oracle باید دقیق باشد.",
        ], sz=14)

        s = self.slide()
        self.title(s, "نمونه کد — Mutant BS-03 + تست Kill")
        self.body(s, [
            "# Mutant (sorting_mutants.py):",
            "for j in range(0, n - i - 2):  # MUTANT",
            "",
            "# تست (test_mutation_killing.py):",
            "def test_m03_bs_killed_wrong_output(...):",
            "    result = bubble_sort_m03(ITEMS_4)  # [d,c,b,a]",
            "    correct = bubble_sort(ITEMS_4)",
            "    assert got != correct",
            "    # [b,c,d,a] != [a,b,c,d] → KILLED ✓",
        ], sz=12)

        s = self.slide()
        self.title(s, "کدام Suite کدام Mutant را می‌کشد؟")
        self.table(s, ["Mutant", "Suite", "نوع fail"], [
            ("BS-01..04", "test_mutation_killing.py", "Exception / assert"),
            ("MS-01,02", "test_mutation_killing.py", "RecursionError"),
            ("MS-03", "—", "EQUIVALENT"),
            ("BN-01,02,04,05", "test_mutation_killing.py", "wrong / timeout"),
            ("BN-03", "—", "EQUIVALENT"),
            ("PC-02", "ACOC (330 tests)", "assert total_after"),
            ("PC-01", "ACOC + assert جدید", "assert total_before"),
        ])

        s = self.slide()
        self.title(s, "سؤال ۱ — کدام معیار پوشش بهتر بود؟")
        self.table(s, ["معیار", "قوت", "ضعف در ctlg"], [
            ("ACOC", "همه ترکیب پارامتر", "PC-01 با oracle ضعیف LIVE شد"),
            ("CFG Node/Edge/Prime", "ساختار و شاخه‌ها", "فضای ورودی کامل نیست"),
            ("Mutation", "کیفیت واقعی تست", "هزینه ساخت mutant"),
        ])
        self.body(s, [
            "پاسخ: Mutation قوی‌ترین معیار کیفیت (91.7%).",
            "CFG برای طراحی تست ساختاری ضروری.",
            "ACOC برای پوشش systematic ورودی.",
            "هر سه مکمل یکدیگرند — 376 تست.",
        ], top=4.0, sz=13)

        s = self.slide()
        self.title(s, "سؤال ۲ — کدام تست‌ها شکست خوردند؟")
        self.body(s, [
            "روی برنامه اصل (Original):",
            "  376/376 passed — هیچ تستی fail نشد ✓",
            "",
            "روی Mutantها (مطلوب است fail شوند):",
            "  11 mutant → KILLED (تست mutation fail)",
            "  2 mutant → EQUIVALENT (قابل kill نیست)",
            "",
            "PC-01 با suite قدیمی LIVE بود —",
            "  تست اصل pass · mutant از assert ضعیف عبور می‌کرد.",
        ], sz=14)

        s = self.slide()
        self.title(s, "سؤال ۳ — بهبود Mutation Score")
        self.body(s, [
            "✓ انجام شد: assert total_before == len(items)",
            "  83.3% → 91.7%",
            "",
            "MS-03, BN-03: EQUIVALENT — از مخرج حذف",
            "  بالاتر از 91.7% بدون mutant جدید ممکن نیست",
            "",
            "پیشنهادهای بیشتر:",
            "  · assert روی meta.sort_algorithm, meta.query",
            "  · boundary: n=12 vs n=13 (switch bubble→merge)",
            "  · property-based tests برای sort invariant",
        ], sz=14)

        s = self.slide()
        self.title(s, "جمع‌بندی و منابع")
        self.image(s, "test_overview_chart.png", 0.35, 1.25, 4.8)
        self.body(s, [
            "UUT: process_catalog + 4 algorithm",
            "ACOC: 330 | CFG: 29 | Mutation: 17",
            "Mutants: 14 AOR | Score: 91.7%",
            "",
            "pytest: cd silver_project",
            "  python -m pytest algorithms/tests/ -v",
            "",
            "HTML: arayeh-kamel.html · graf-ha.html",
            "گراف: cfg_graphs/*.html",
        ], top=1.3, sz=12)

        s = self.slide()
        deco = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.2), Inches(6), Inches(2.8))
        deco.fill.solid()
        deco.fill.fore_color.rgb = C_WHITE
        deco.line.color.rgb = C_PURPLE_SOFT
        tb = s.shapes.add_textbox(Inches(2.2), Inches(2.5), Inches(5.6), Inches(2.2))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, (t, sz, c) in enumerate([
            ("سپاس از توجه شما", 28, C_PINK),
            ("سؤالی هست؟", 18, C_MUTED),
            ("376 Tests · 91.7% Mutation Score", 14, C_PURPLE),
        ]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = t
            p.font.size = Pt(sz)
            p.font.color.rgb = c
            p.font.bold = i == 0
            p.alignment = PP_ALIGN.CENTER

        OUT.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(OUT))
        self.prs.save(str(OUT2))
        print(f"Saved: {OUT} ({self.n} slides)")
        print(f"Saved: {OUT2}")


if __name__ == "__main__":
    Deck().build_all()
